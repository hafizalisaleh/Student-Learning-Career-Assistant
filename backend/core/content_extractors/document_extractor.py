"""
Document content extractors for various file formats
"""
import os
import base64
import json
import requests
from typing import Optional
from pathlib import Path
import PyPDF2
from docx import Document as DocxDocument
import pptx
import pandas as pd
from PIL import Image
from config.settings import settings

class DocumentExtractor:
    """Extract content from various document formats"""
    
    def __init__(self):
        self.ocr_api_key = settings.OCR_API_KEY
        self.ocr_api_secret = settings.OCR_API_SECRET
        self.ocr_api_url = settings.OCR_API_URL
    
    def extract_from_pdf(self, file_path: str) -> Optional[str]:
        """
        Extract text from PDF file
        
        Args:
            file_path: Path to PDF file
            
        Returns:
            Extracted text
        """
        try:
            text = ""
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
            return text
        except Exception as e:
            print(f"Error reading PDF: {e}")
            return None
    
    def extract_from_docx(self, file_path: str) -> Optional[str]:
        """
        Extract text from Word document
        
        Args:
            file_path: Path to DOCX file
            
        Returns:
            Extracted text
        """
        try:
            doc = DocxDocument(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            print(f"Error reading Word document: {e}")
            return None
    
    def extract_from_pptx(self, file_path: str) -> Optional[str]:
        """
        Extract text from PowerPoint presentation
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            Extracted text
        """
        try:
            presentation = pptx.Presentation(file_path)
            text = ""
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            print(f"Error reading PowerPoint: {e}")
            return None
    
    def extract_from_excel(self, file_path: str) -> Optional[str]:
        """
        Extract text from Excel file
        
        Args:
            file_path: Path to Excel file
            
        Returns:
            Extracted text
        """
        try:
            text = ""
            excel_file = pd.ExcelFile(file_path)
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                text += f"Sheet: {sheet_name}\n"
                text += df.to_string() + "\n\n"
            return text
        except Exception as e:
            print(f"Error reading Excel: {e}")
            return None
    
    def extract_from_txt(self, file_path: str) -> Optional[str]:
        """
        Extract text from text file
        
        Args:
            file_path: Path to text file
            
        Returns:
            Extracted text
        """
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except Exception as e:
            print(f"Error reading text file: {e}")
            return None
    
    def ocr_image(self, image_path: str) -> Optional[str]:
        """
        Perform OCR on image using OCR API
        
        Args:
            image_path: Path to image file
            
        Returns:
            Extracted text
        """
        try:
            # Convert image to Base64
            with open(image_path, "rb") as img_file:
                image_b64 = base64.b64encode(img_file.read()).decode("utf-8")
            
            # Prepare payload
            payload = {
                "api_key": self.ocr_api_key,
                "api_secret": self.ocr_api_secret,
                "image": image_b64
            }
            
            # Send POST request
            headers = {"Content-Type": "application/json; charset=utf-8"}
            response = requests.post(self.ocr_api_url, headers=headers, data=json.dumps(payload))
            
            if response.status_code == 200:
                return response.text
            else:
                print(f"OCR Error: {response.status_code}, {response.text}")
                return None
                
        except Exception as e:
            print(f"Error during OCR: {e}")
            return None
    
    def extract_from_image(self, file_path: str) -> Optional[str]:
        """
        Returns special marker for image files to be processed by Gemini Vision
        Instead of OCR, we'll pass the image directly to Gemini
        
        Args:
            file_path: Path to image file
            
        Returns:
            Special marker with image path for Gemini Vision processing
        """
        # Return special marker that will be detected by Gemini client
        return f"__GEMINI_IMAGE__{file_path}__"
    
    def extract_text(self, file_path: str) -> Optional[str]:
        """
        Extract text from file based on extension
        
        Args:
            file_path: Path to file
            
        Returns:
            Extracted text
        """
        ext = Path(file_path).suffix.lower()
        
        extractors = {
            '.pdf': self.extract_from_pdf,
            '.docx': self.extract_from_docx,
            '.doc': self.extract_from_docx,
            '.pptx': self.extract_from_pptx,
            '.ppt': self.extract_from_pptx,
            '.xlsx': self.extract_from_excel,
            '.xls': self.extract_from_excel,
            '.csv': self.extract_from_excel,
            '.txt': self.extract_from_txt,
            '.md': self.extract_from_txt,
            '.jpg': self.extract_from_image,
            '.jpeg': self.extract_from_image,
            '.png': self.extract_from_image,
            '.bmp': self.extract_from_image,
            '.tiff': self.extract_from_image,
            '.tif': self.extract_from_image,
        }
        
        extractor = extractors.get(ext)
        if extractor:
            return extractor(file_path)
        else:
            print(f"Unsupported file type: {ext}")
            return None
